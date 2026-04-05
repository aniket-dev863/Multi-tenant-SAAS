from pptx import Presentation
from pptx.util import Pt

def add_bullet(tf, bold_text, normal_text="", level=0, is_first=False):
    """Helper function to create rich, multi-level bullets with bold headers."""
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.level = level
    
    if bold_text:
        run = p.add_run()
        run.text = bold_text
        run.font.bold = True
    if normal_text:
        run = p.add_run()
        run.text = normal_text

def create_rich_presentation():
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # Slide 1: Title
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "StreamPulse Pharma"
    slide1.placeholders[1].text = "Cloud-Native Multi-Tenant SaaS Platform\n\nPresented By: Aniket Vyavahare\nSoftware Engineering Evaluation"

    # Slide 2: Problem
    slide2 = prs.slides.add_slide(bullet_layout)
    slide2.shapes.title.text = "The Legacy Pharmacy Challenge"
    tf2 = slide2.placeholders[1].text_frame
    add_bullet(tf2, "The Monolith Bottleneck:", " Traditional software lacks real-time branch syncing.", level=0, is_first=True)
    add_bullet(tf2, "Result:", " Inventory mismatches across different city branches.", level=1)
    add_bullet(tf2, "Hardware Dependency:", " Local server crashes cause catastrophic data loss.", level=0)
    add_bullet(tf2, "Result:", " Loss of crucial sales and compliance records.", level=1)
    add_bullet(tf2, "Security Risks:", " Single-tenant setups lack enterprise-grade failover.", level=0)

    # Slide 3: Solution
    slide3 = prs.slides.add_slide(bullet_layout)
    slide3.shapes.title.text = "The Proposed SaaS Solution"
    tf3 = slide3.placeholders[1].text_frame
    add_bullet(tf3, "Cloud-First Architecture:", " Migrating to a highly available AWS environment.", level=0, is_first=True)
    add_bullet(tf3, "Multi-Tenant Design:", " One robust application serving multiple independent chains.", level=0)
    add_bullet(tf3, "Benefit:", " Drastically reduces infrastructure costs while scaling instantly.", level=1)
    add_bullet(tf3, "Hybrid Data Model:", " Utilizing both SQL and NoSQL databases.", level=0)

    # Slide 4: Multi-Tenant Strategy
    slide4 = prs.slides.add_slide(bullet_layout)
    slide4.shapes.title.text = "Multi-Tenant Architecture Strategy"
    tf4 = slide4.placeholders[1].text_frame
    add_bullet(tf4, "Discriminator Column Isolation:", " Every record is tied to a specific tenant ID.", level=0, is_first=True)
    add_bullet(tf4, "Implementation:", " A shared 'pharmacies' table links all data.", level=1)
    add_bullet(tf4, "Session-Level Security:", " Backend strictly scopes queries to the logged-in tenant.", level=0)
    add_bullet(tf4, "Guarantee:", " Cross-organization data leaks are mathematically impossible.", level=1)
    add_bullet(tf4, "Resource Efficiency:", " Maximum compute utilization.", level=0)

    # Slide 5: Hybrid DB
    slide5 = prs.slides.add_slide(bullet_layout)
    slide5.shapes.title.text = "Hybrid Database Implementation"
    tf5 = slide5.placeholders[1].text_frame
    add_bullet(tf5, "Amazon RDS (PostgreSQL):", " Strict, ACID-compliant data.", level=0, is_first=True)
    add_bullet(tf5, "Use Case:", " Users, Sales, Inventory Batches, and Financial Ledgers.", level=1)
    add_bullet(tf5, "Amazon DocumentDB (MongoDB):", " Schema-less, flexible data.", level=0)
    add_bullet(tf5, "Use Case:", " Dynamic medicine side-effects, descriptions, and supplier notes.", level=1)
    add_bullet(tf5, "Why?", " Polyglot persistence ensures the right tool for the right data.", level=0)

    # Slide 6: Architecture Diagram
    slide6 = prs.slides.add_slide(bullet_layout)
    slide6.shapes.title.text = "Enterprise Cloud Architecture"
    tf6 = slide6.placeholders[1].text_frame
    add_bullet(tf6, "[ ⚠️ DRAG AND DROP YOUR AWS DIAGRAM IMAGE HERE ]", level=0, is_first=True)
    add_bullet(tf6, "Key Components:", level=0)
    add_bullet(tf6, "Application Load Balancer (ALB) & Auto Scaling Group", level=1)
    add_bullet(tf6, "Multi-AZ Database Replication for Failover", level=1)

    # Slide 7: Security Posture
    slide7 = prs.slides.add_slide(bullet_layout)
    slide7.shapes.title.text = "Security & Compliance Posture"
    tf7 = slide7.placeholders[1].text_frame
    add_bullet(tf7, "Network Isolation (VPC):", level=0, is_first=True)
    add_bullet(tf7, "Public Subnet:", " Houses the EC2 Web Servers.", level=1)
    add_bullet(tf7, "Private Subnet:", " Secures the RDS and DocumentDB instances from the internet.", level=1)
    add_bullet(tf7, "Secret Management:", " Dynamic credential injection via .env / AWS Secrets Manager.", level=0)
    add_bullet(tf7, "Application Layer:", " Strict Role-Based Access Control (Admin vs. Pharmacist).", level=0)

    # Slide 8: Concurrency
    slide8 = prs.slides.add_slide(bullet_layout)
    slide8.shapes.title.text = "Handling Concurrency & Scale"
    tf8 = slide8.placeholders[1].text_frame
    add_bullet(tf8, "Asynchronous Processing:", " Replaced development server with Gunicorn.", level=0, is_first=True)
    add_bullet(tf8, "Workers:", " Multiple worker processes handle parallel requests.", level=1)
    add_bullet(tf8, "Elastic Scalability:", " Handles hundreds of simultaneous POS checkouts.", level=0)
    add_bullet(tf8, "Stateless Architecture:", " Signed sessions allow seamless routing across servers.", level=0)

    # Slide 9: Future Scope
    slide9 = prs.slides.add_slide(bullet_layout)
    slide9.shapes.title.text = "Future Scope: Event-Driven Automation"
    tf9 = slide9.placeholders[1].text_frame
    add_bullet(tf9, "Moving Beyond Dashboards:", " Implementing active, event-driven alerts.", level=0, is_first=True)
    add_bullet(tf9, "Database Triggers:", " Low inventory automatically invokes AWS Lambda.", level=0)
    add_bullet(tf9, "Instant Notifications:", " Lambda triggers Amazon SNS.", level=1)
    add_bullet(tf9, "Action:", " Immediately SMS/Emails the manager to reorder stock.", level=1)

    # Slide 10: Conclusion
    slide10 = prs.slides.add_slide(bullet_layout)
    slide10.shapes.title.text = "Current Progress & Roadmap"
    tf10 = slide10.placeholders[1].text_frame
    add_bullet(tf10, "Phase 1 Complete (50% Milestone):", level=0, is_first=True)
    add_bullet(tf10, "Discriminator column logic and Hybrid DB successfully implemented locally.", level=1)
    add_bullet(tf10, "Phase 2 Roadmap:", level=0)
    add_bullet(tf10, "Full deployment to AWS EC2, RDS, and configuring Lambda triggers.", level=1)
    add_bullet(tf10, "Thank you.", " Open for questions regarding the architecture.", level=0)

    prs.save('StreamPulse_Rich_Presentation.pptx')
    print("Rich presentation successfully created: StreamPulse_Rich_Presentation.pptx")

if __name__ == '__main__':
    create_rich_presentation()